from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os, textwrap

ROOT = os.getcwd()
PAGE_PATH = os.path.join(ROOT, "app", "page.tsx")
OUT_PATH = os.path.join(ROOT, "AI2_Testing_Report.pptx")

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

def add_title(title, subtitle=None):
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = title
    if subtitle and len(s.placeholders) > 1:
        s.placeholders[1].text = subtitle

def add_bullets(title, bullets):
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = title
    tf = s.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0

def add_text_slide(title, text):
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = title
    tb = s.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12.3), Inches(5.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.font.size = Pt(10)
    p.text = text

# Slides content
add_title("AI 2 PROJECT — Testing & Debugging Report", "Realistic Earth Scene + Image Upload — Test plan, findings, fixes")
add_bullets("Agenda", [
    "Project summary",
    "Architecture & key files",
    "Current code issues (from app/page.tsx)",
    "Test strategy & matrix",
    "Representative test cases",
    "Debugging workflow & fixes",
    "CI / automation & next steps"
])

add_bullets("Project summary", [
    "Next.js app integrating a cinematic Three.js globe (react-three-fiber).",
    "Interactive features: drag-to-rotate, smooth zoom, clouds, atmosphere, rings, hubs.",
    "Image upload + backend APIs: /api/environment, /api/value, /api/summary."
])

add_bullets("Architecture & key files", [
    "app/page.tsx — main page wiring scene, upload, UI",
    "components/RealisticEarthScene.tsx — 3D scene",
    "components/ImageUpload, ResultsPanel, ChatBox, Navigation",
    "public/textures/ — earth_day, earth_night, bump, specular, clouds",
    "scripts/ — automation and PPT generator"
])

# Include current app/page.tsx excerpt
page_text = ""
if os.path.exists(PAGE_PATH):
    with open(PAGE_PATH, "r", encoding="utf-8", errors="ignore") as f:
        page_text = f.read()
    excerpt = page_text[:3500]
else:
    excerpt = "// app/page.tsx not found in project. Place the file at app/page.tsx to include its content."

add_text_slide("Current app/page.tsx (excerpt)", excerpt)

# Current code issues & fixes
issues = [
    "Detected error in app/page.tsx imports:",
    "1) Invalid import syntax: `import - from '@/components/ChatBox'` (causes parse error).",
    "2) Dynamic import uses RealisticEarth2050 instead of RealisticEarthScene.",
    "Suggested fixes:",
    "A) Replace invalid import with: `import ChatBox from '@/components/ChatBox'`",
    "B) Replace dynamic import line with: `const Scene = dynamic(() => import('@/components/RealisticEarthScene'), { ssr: false })`"
]
add_bullets("Current code issues & suggested fixes", issues)

add_bullets("Test strategy", [
    "Unit tests: utilities (latLonToVector3), small components, error fallback behavior.",
    "Integration: ImageUpload -> API -> ResultsPanel (mock APIs with MSW).",
    "E2E: Playwright flow for upload -> results -> download with traces/screenshots.",
    "Performance: FPS/memory profiling and shader fallbacks for low-end GPUs."
])

add_bullets("Test matrix (priority)", [
    "High: ImageUpload, API endpoints, RealisticEarthScene loading & interaction",
    "Medium: OrbitControls, cloud/atmosphere shader behaviors",
    "Low: Visual parity of high-res textures and advanced shader improvements"
])

add_bullets("Representative test cases", [
    "Unit: latLonToVector3 returns expected vectors for SF (37.7749,-122.4194) and London (51.5074,-0.1278).",
    "Unit: Scene mounts without throwing when textures missing (graceful fallback).",
    "Integration: POST /api/environment returns 200 and ResultsPanel displays values.",
    "E2E: User uploads sample image with GPS -> map pin placed -> ResultsPanel visible -> download completes.",
    "Performance: Maintain >=30 FPS during 2-minute interactive session on target hardware."
])

add_bullets("Debugging workflow", [
    "1) Reproduce reliably with sample files and exact steps.",
    "2) Isolate: disable heavy postprocessing or complex shaders to narrow issues.",
    "3) Add targeted console logs and capture browser console/network traces.",
    "4) Use React DevTools, Chrome GPU profiler, and VSCode attach debugger.",
    "5) For GL/shader issues: fallback to basic material and validate textures (format, colorSpace)."
])

add_bullets("CI / automation recommendations", [
    "Run: typecheck (tsc --noEmit), lint (eslint), unit tests (jest), e2e (playwright).",
    "Use MSW to mock APIs in tests; upload Playwright artifacts on failure.",
    "Fail PRs on typecheck/test failures and publish coverage/report artifacts."
])

add_bullets("How to run locally (short)", [
    'cd "C:\\Users\\Mohammed Hamdan\\AI 2 PROJECT"',
    "pip install python-pptx (for this generator)",
    "npm install",
    "npm run dev  (or npx next dev)",
    "Unit tests: npm run test   E2E: npx playwright test"
])

add_bullets("Performance & optimization checklist", [
    "Convert textures to KTX2/Basis and use KTX2 loader (smaller GPU uploads).",
    "Wrap heavy assets in <Suspense> and add <Preload all />.",
    "Provide low-power mode: reduce sphere segments, disable bloom selectively.",
    "Profile and capture FPS traces for key target devices."
])

add_bullets("Evidence & artifacts (placeholders)", [
    "Include: VSCode error screenshot, terminal test log (logs/test_results.txt), Playwright traces/screenshots.",
    "Add before/after screenshots of scene and code-diff screenshots for fixes."
])

add_bullets("Recommended fixes & PR checklist", [
    "Fix imports in app/page.tsx (see suggested fixes slide).",
    "Add unit tests for utilities and a mount test for RealisticEarthScene (mock textures).",
    "Add MSW mocked API integration tests and Playwright E2E for critical flows.",
    "Add CI pipeline: typecheck, lint, test, e2e, publish artifacts."
])

# Append full path footer slide
add_bullets("Generated", [
    f"Project path: {ROOT}",
    f"Included file excerpt: {PAGE_PATH if os.path.exists(PAGE_PATH) else 'not found'}",
    "Script: scripts/generate_ai2_ppt.py"
])

# Save PPT
prs.save(OUT_PATH)
print("Saved PPT:", OUT_PATH)
